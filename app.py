import streamlit as st
import io
import os
import json
from pptx import Presentation
import qrcode
from PIL import Image
from datetime import datetime, date

# ── PAGE CONFIG ───────────────────────────────────────────────
st.set_page_config(
    page_title="DR Factsheet Generator",
    page_icon="📊",
    layout="wide",
)

# ── PASSWORD GATE ─────────────────────────────────────────────
APP_PASSWORD = "Password1234"  # ← change this

if "authenticated" not in st.session_state:
    st.session_state.authenticated = False

if not st.session_state.authenticated:
    st.markdown("""
    <div style="max-width:380px;margin:80px auto;">
        <div style="background:#141720;border:1px solid #252A3A;border-radius:14px;
                    padding:36px 32px;text-align:center;">
            <div style="font-size:36px;margin-bottom:8px;">🔒</div>
            <div style="font-size:18px;font-weight:600;color:#E8ECF4;margin-bottom:4px;">
                DR Factsheet Generator</div>
            <div style="font-size:12px;color:#5A637A;margin-bottom:24px;">
                กรุณาใส่รหัสผ่านเพื่อเข้าใช้งาน</div>
        </div>
    </div>
    """, unsafe_allow_html=True)
    with st.form("login_form"):
        _, mid, _ = st.columns([1, 2, 1])
        with mid:
            password = st.text_input("รหัสผ่าน", type="password", placeholder="Enter password")
            login = st.form_submit_button("เข้าสู่ระบบ", use_container_width=True, type="primary")
    if login:
        if password == APP_PASSWORD:
            st.session_state.authenticated = True
            st.rerun()
        else:
            _, mid, _ = st.columns([1, 2, 1])
            with mid:
                st.error("❌ รหัสผ่านไม่ถูกต้อง")
    st.stop()

# ── FIXED VALUES ──────────────────────────────────────────────
FIXED = {
    "exchange"     : "ตลาดหลักทรัพย์แห่งประเทศไทย (SET)",
    "depositary"   : "ธนาคารกรุงไทย จำกัด (มหาชน)",
    "offering_type": "Direct Listing",
    "price_info"   : "เป็นไปตามกลไกราคาของตลาดในเวลาที่เสนอขาย โดยอ้างอิงจากราคาหลักทรัพย์อ้างอิงต่างประเทศ อัตราแลกเปลี่ยน อัตราส่วนอ้างอิง และค่าธรรมเนียมอื่น ๆ",
    "ktb_contact"  : "02-208-3748, 02-208-4669",
}

# ── EXCHANGE DROPDOWN OPTIONS ─────────────────────────────────
# Format: "Full Exchange Name" : "SHORT"
EXCHANGES = {
    "Hong Kong Stock Exchange"      : "HKEX",
    "Nasdaq Stock Exchange"         : "NASDAQ",
    "New York Stock Exchange"       : "NYSE",
    "Tokyo Stock Exchange"          : "TSE",
    "Shanghai Stock Exchange"       : "SSE",
    "Shenzhen Stock Exchange"       : "SZSE",
    "Euronext Paris"                : "EPA",
    "London Stock Exchange"         : "LSE",
    "Singapore Exchange"            : "SGX",
    "Korea Stock Exchange"          : "KRX",
    "Taiwan Stock Exchange"         : "TWSE",
    "Australian Securities Exchange": "ASX",
}

# ── TEMPLATE FILE PATH (persisted on disk next to app.py) ─────
TEMPLATE_PATH = os.path.join(os.path.dirname(os.path.abspath(__file__)), "template.pptx")

def load_template_from_disk():
    """Load template bytes from disk if file exists."""
    if os.path.exists(TEMPLATE_PATH):
        with open(TEMPLATE_PATH, "rb") as f:
            return f.read()
    return None

def save_template_to_disk(data: bytes):
    """Save template bytes to disk."""
    with open(TEMPLATE_PATH, "wb") as f:
        f.write(data)

# ── HISTORY FILE PATH (persisted on disk next to app.py) ──────
HISTORY_PATH = os.path.join(os.path.dirname(os.path.abspath(__file__)), "history.json")

def load_history_from_disk():
    """Load history list from disk. Returns [] if missing or corrupt."""
    if not os.path.exists(HISTORY_PATH):
        return []
    try:
        with open(HISTORY_PATH, "r", encoding="utf-8") as f:
            data = json.load(f)
        # Rehydrate _trading_date strings back to date objects
        for entry in data:
            if isinstance(entry.get("_trading_date"), str):
                try:
                    entry["_trading_date"] = date.fromisoformat(entry["_trading_date"])
                except ValueError:
                    entry["_trading_date"] = date.today()
        return data if isinstance(data, list) else []
    except Exception:
        return []

def save_history_to_disk(history: list):
    """Save history list to disk as JSON."""
    serialisable = []
    for entry in history:
        row = dict(entry)
        if isinstance(row.get("_trading_date"), date):
            row["_trading_date"] = row["_trading_date"].isoformat()
        serialisable.append(row)
    with open(HISTORY_PATH, "w", encoding="utf-8") as f:
        json.dump(serialisable, f, ensure_ascii=False, indent=2)

# ── THAI DATE FORMATTER ───────────────────────────────────────
THAI_MONTHS = ["ม.ค.","ก.พ.","มี.ค.","เม.ย.","พ.ค.","มิ.ย.",
               "ก.ค.","ส.ค.","ก.ย.","ต.ค.","พ.ย.","ธ.ค."]

def to_thai_date(d: date) -> str:
    thai_year = d.year - 1957
    return f"{d.day} {THAI_MONTHS[d.month - 1]} {thai_year}"

# ── SESSION STATE ─────────────────────────────────────────────
if "form_data"       not in st.session_state: st.session_state.form_data       = None
if "history"         not in st.session_state: st.session_state.history         = load_history_from_disk()
if "template_bytes"  not in st.session_state:
    # Try to load from disk on first run
    st.session_state.template_bytes = load_template_from_disk()
if "edit_index"      not in st.session_state: st.session_state.edit_index      = None
if "prefill"         not in st.session_state: st.session_state.prefill         = {}
if "selected_exchange" not in st.session_state: st.session_state.selected_exchange = list(EXCHANGES.keys())[0]

# ── STYLES ────────────────────────────────────────────────────
st.markdown("""
<style>
    @import url('https://fonts.googleapis.com/css2?family=IBM+Plex+Sans+Thai:wght@300;400;500;600&display=swap');
    html, body, [class*="css"] { font-family: 'IBM Plex Sans Thai', sans-serif; }
    .block-container { padding: 2rem 2.5rem; max-width: 1400px; }
    .app-header {
        background: linear-gradient(135deg,#141720 0%,#1C2030 100%);
        border: 1px solid #252A3A; border-radius: 12px;
        padding: 20px 28px; margin-bottom: 24px;
        display: flex; align-items: center; gap: 16px;
    }
    .header-badge {
        background: #3B6FFF; color: white; font-weight: 700; font-size: 18px;
        width: 44px; height: 44px; border-radius: 10px;
        display: flex; align-items: center; justify-content: center;
    }
    .section-label {
        font-size: 11px; font-weight: 600; letter-spacing: 0.1em;
        text-transform: uppercase; color: #5A637A; margin: 16px 0 8px;
    }
    .fixed-grid { display:grid; grid-template-columns:repeat(3,1fr); gap:10px; margin-bottom:16px; }
    .fixed-card { background:#141720; border:1px solid #252A3A; border-radius:8px; padding:10px 14px; }
    .fixed-label { font-size:10px; color:#5A637A; text-transform:uppercase; letter-spacing:.06em; }
    .fixed-value { font-size:12px; color:#8892A4; margin-top:3px; }
    .gen-card { background:#0D1A12; border:1px solid #1A3A25; border-radius:8px;
                padding:10px 14px; margin-bottom:6px; }
    .gen-label { font-size:10px; color:#00D4AA; text-transform:uppercase; letter-spacing:.06em; }
    .gen-value { font-size:12px; color:#B0C4B8; margin-top:3px; line-height:1.5; }
    .hist-card {
        background:#141720; border:1px solid #252A3A; border-radius:8px;
        padding:12px 14px; margin-bottom:4px;
    }
    .hist-ticker {
        display:inline-block;
        background:rgba(59,111,255,.15); color:#7DA4FF;
        border:1px solid rgba(59,111,255,.25); border-radius:4px;
        padding:2px 10px; font-size:12px; font-weight:600;
        font-family:'IBM Plex Mono',monospace; margin-bottom:4px;
    }
    .hist-meta { font-size:11px; color:#5A637A; }
    .stButton>button { font-family:'IBM Plex Sans Thai',sans-serif !important;
                       border-radius:8px !important; font-weight:500 !important; }
    div[data-testid="stForm"] { background:#141720; border:1px solid #252A3A;
                                 border-radius:12px; padding:20px; }
    .template-status-ok {
        background:#0D1A12; border:1px solid #1A3A25; border-radius:8px;
        padding:8px 14px; font-size:12px; color:#00D4AA; margin-bottom:8px;
    }
    .template-status-missing {
        background:#1A1212; border:1px solid #3A1A1A; border-radius:8px;
        padding:8px 14px; font-size:12px; color:#FF6B6B; margin-bottom:8px;
    }
    hr { border-color:#252A3A !important; }
</style>
""", unsafe_allow_html=True)

# ── HEADER ────────────────────────────────────────────────────
header_col, logout_col = st.columns([6, 1])
with header_col:
    st.markdown("""
    <div class="app-header">
        <div class="header-badge">DR</div>
        <div>
            <div style="font-size:18px;font-weight:600;color:#E8ECF4;">DR Factsheet Generator</div>
            <div style="font-size:12px;color:#5A637A;margin-top:2px;">กรอกข้อมูลเพื่อสร้าง Factsheet อัตโนมัติ</div>
        </div>
    </div>
    """, unsafe_allow_html=True)
with logout_col:
    st.markdown("<div style='height:18px'></div>", unsafe_allow_html=True)
    if st.button("🔓 ออกจากระบบ", use_container_width=True):
        st.session_state.authenticated = False
        st.session_state.form_data     = None
        st.rerun()

# ── HELPERS ───────────────────────────────────────────────────
def build_data(ticker, company_name, stock_code, exchange_full,
               total_units, ratio, trading_date_obj, filing_url, exchange_short):
    thai_date  = to_thai_date(trading_date_obj)
    fmt_units  = f"{int(total_units):,} หน่วย"
    fmt_ratio  = f"1 : {int(ratio):,}"
    full_thai  = (f"ตราสารแสดงสิทธิในหลักทรัพย์ต่างประเทศของบริษัท {company_name} "
                  f"ออกโดยธนาคารกรุงไทย จำกัด (มหาชน)")
    full_eng   = (f"Depositary receipt on {company_name} "
                  f"issued by Krungthai Bank Public Company Limited")
    underlying = f"{company_name} ({stock_code})"
    return {
        "ticker"             : ticker,
        "full_name_thai"     : full_thai,
        "full_name_eng"      : full_eng,
        "exchange"           : FIXED["exchange"],
        "underlying_stock"   : underlying,
        "underlying_exchange": exchange_full,
        "depositary"         : FIXED["depositary"],
        "offering_type"      : FIXED["offering_type"],
        "total_units"        : fmt_units,
        "ratio"              : fmt_ratio,
        "first_trading_date" : thai_date,
        "price_info"         : FIXED["price_info"],
        "ktb_contact"        : FIXED["ktb_contact"],
        "filing_url"         : filing_url,
        "foreign_exchange"   : exchange_short,
        # raw inputs stored for re-editing
        "_ticker"        : ticker,
        "_company_name"  : company_name,
        "_stock_code"    : stock_code,
        "_exchange_full" : exchange_full,
        "_exchange_short": exchange_short,
        "_total_units"   : int(total_units),
        "_ratio"         : int(ratio),
        "_trading_date"  : trading_date_obj,
        "_filing_url"    : filing_url,
        "_created"       : datetime.now().strftime("%d/%m/%y %H:%M"),
    }

def replace_in_paragraph(paragraph, placeholder, value):
    full_text = "".join(run.text for run in paragraph.runs)
    if placeholder not in full_text:
        return
    new_text = full_text.replace(placeholder, value)
    if paragraph.runs:
        paragraph.runs[0].text = new_text
        for run in paragraph.runs[1:]:
            run.text = ""

def replace_all_placeholders(prs, data):
    pmap = {
        "{{ticker}}"             : data["ticker"],
        "{{full_name_thai}}"     : data["full_name_thai"],
        "{{full_name_eng}}"      : data["full_name_eng"],
        "{{exchange}}"           : data["exchange"],
        "{{underlying_stock}}"   : data["underlying_stock"],
        "{{underlying_exchange}}": data["underlying_exchange"],
        "{{depositary}}"         : data["depositary"],
        "{{offering_type}}"      : data["offering_type"],
        "{{total_units}}"        : data["total_units"],
        "{{ratio}}"              : data["ratio"],
        "{{first_trading_date}}" : data["first_trading_date"],
        "{{price_info}}"         : data["price_info"],
        "{{ktb_contact}}"        : data["ktb_contact"],
        "{{filing_url}}"         : data["filing_url"],
        "{{foreign_exchange}}"   : data["foreign_exchange"],
    }
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for ph, val in pmap.items():
                        replace_in_paragraph(para, ph, str(val) if val else "")
            if shape.shape_type == 19:  # table
                for row in shape.table.rows:
                    for cell in row.cells:
                        for para in cell.text_frame.paragraphs:
                            for ph, val in pmap.items():
                                replace_in_paragraph(para, ph, str(val) if val else "")

def insert_qr_code(prs, url):
    if not url:
        return
    qr = qrcode.QRCode(box_size=10, border=2)
    qr.add_data(url)
    qr.make(fit=True)
    img_bytes = io.BytesIO()
    qr.make_image(fill_color="black", back_color="white").save(img_bytes, format="PNG")
    img_bytes.seek(0)
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip() == "{{qr_code}}":
                left, top, w, h = shape.left, shape.top, shape.width, shape.height
                shape._element.getparent().remove(shape._element)
                slide.shapes.add_picture(img_bytes, left, top, w, h)
                img_bytes.seek(0)

def generate_pptx(template_bytes, data):
    prs = Presentation(io.BytesIO(template_bytes))
    replace_all_placeholders(prs, data)
    if data.get("filing_url"):
        insert_qr_code(prs, data["filing_url"])
    output = io.BytesIO()
    prs.save(output)
    return output.getvalue()

# ══════════════════════════════════════════════════════════════
#  THREE-COLUMN LAYOUT
# ══════════════════════════════════════════════════════════════
form_col, preview_col, history_col = st.columns([3, 2, 2], gap="large")

# ════════════════
#  LEFT — FORM
# ════════════════
with form_col:

    # ── TEMPLATE SECTION ──────────────────────────────────────
    with st.expander("📎 Template (.pptx)", expanded=st.session_state.template_bytes is None):

        if st.session_state.template_bytes:
            st.markdown('<div class="template-status-ok">✅ Template พร้อมใช้งาน (บันทึกไว้ในระบบแล้ว)</div>',
                        unsafe_allow_html=True)
            # Download current template
            st.download_button(
                label="⬇️ ดาวน์โหลด Template ปัจจุบัน",
                data=st.session_state.template_bytes,
                file_name="template.pptx",
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                use_container_width=True,
            )
            st.caption("อัปโหลดไฟล์ใหม่ด้านล่างเพื่อแทนที่ template ปัจจุบัน")
        else:
            st.markdown('<div class="template-status-missing">⚠️ ยังไม่มี Template — กรุณาอัปโหลด .pptx</div>',
                        unsafe_allow_html=True)

        uploaded = st.file_uploader("อัปโหลด template ใหม่", type=["pptx"], label_visibility="collapsed")
        if uploaded:
            new_bytes = uploaded.read()
            st.session_state.template_bytes = new_bytes
            save_template_to_disk(new_bytes)          # ← persist to disk
            st.success(f"✅ บันทึก {uploaded.name} เรียบร้อย — จะใช้งาน template นี้จนกว่าจะอัปโหลดใหม่")
            st.rerun()

    st.markdown("---")
    st.markdown('<div class="section-label">🔒 ข้อมูลคงที่</div>', unsafe_allow_html=True)
    st.markdown(f"""
    <div class="fixed-grid">
        <div class="fixed-card"><div class="fixed-label">ตลาดรอง</div>
            <div class="fixed-value">{FIXED['exchange']}</div></div>
        <div class="fixed-card"><div class="fixed-label">ผู้ออกตราสาร</div>
            <div class="fixed-value">{FIXED['depositary']}</div></div>
        <div class="fixed-card"><div class="fixed-label">รูปแบบ</div>
            <div class="fixed-value">{FIXED['offering_type']}</div></div>
        <div class="fixed-card" style="grid-column:span 2;"><div class="fixed-label">ราคาตราสาร</div>
            <div class="fixed-value">{FIXED['price_info']}</div></div>
        <div class="fixed-card"><div class="fixed-label">เบอร์ KTB</div>
            <div class="fixed-value">{FIXED['ktb_contact']}</div></div>
    </div>
    """, unsafe_allow_html=True)

    st.markdown("---")

    pf      = st.session_state.prefill
    editing = st.session_state.edit_index is not None

    if editing:
        st.info(f"✏️ แก้ไขรายการ: **{pf.get('_ticker', '')}**")

    st.markdown('<div class="section-label">✏️ กรอกข้อมูล DR</div>', unsafe_allow_html=True)

    # ── EXCHANGE SELECTOR (outside form so short name updates live) ──
    exchange_options = list(EXCHANGES.keys())
    saved_exch       = pf.get("_exchange_full", st.session_state.selected_exchange)
    default_idx      = exchange_options.index(saved_exch) if saved_exch in exchange_options else 0

    # Determine the correct short name to display:
    # Priority: 1) user manually typed something  2) prefill from edit  3) auto from map
    def _init_short_name():
        """Called once when exchange changes or prefill changes."""
        if pf.get("_exchange_full") == st.session_state.get("_cur_exchange_full") and pf.get("_exchange_short"):
            return pf["_exchange_short"]
        return EXCHANGES.get(st.session_state.get("_cur_exchange_full", saved_exch), "")

    # Detect if the selectbox value changed — if so, reset the short name
    prev_exchange = st.session_state.get("_cur_exchange_full")

    col3, col4 = st.columns(2)
    with col3:
        exchange_full = st.selectbox(
            "ตลาดจดทะเบียน *",
            options=exchange_options,
            index=default_idx,
            key="exchange_selectbox",
        )

    # When exchange changes, update the displayed short name automatically
    if prev_exchange != exchange_full:
        st.session_state["_cur_exchange_full"] = exchange_full
        # Use prefill short if we're editing and this exchange matches
        if pf.get("_exchange_full") == exchange_full and pf.get("_exchange_short"):
            st.session_state["_short_name_display"] = pf["_exchange_short"]
        else:
            st.session_state["_short_name_display"] = EXCHANGES.get(exchange_full, "")

    # First-time init
    if "_short_name_display" not in st.session_state:
        if pf.get("_exchange_full") == exchange_full and pf.get("_exchange_short"):
            st.session_state["_short_name_display"] = pf["_exchange_short"]
        else:
            st.session_state["_short_name_display"] = EXCHANGES.get(exchange_full, "")

    with col4:
        # No key= here so value= always wins; we capture the result manually
        exchange_short_input = st.text_input(
            "ชื่อย่อตลาด *",
            value=st.session_state["_short_name_display"],
            help="กำหนดอัตโนมัติจากตลาดที่เลือก แต่สามารถแก้ไขได้",
        )
        # Keep session state in sync with whatever user typed
        st.session_state["_short_name_display"] = exchange_short_input

    with st.form("dr_form", clear_on_submit=False):

        col1, col2 = st.columns(2)
        with col1:
            ticker = st.text_input("ชื่อย่อ (Ticker) *",
                                   value=pf.get("_ticker", ""),
                                   placeholder="เช่น SUNNY80")
        with col2:
            stock_code = st.text_input("รหัสหลักทรัพย์อ้างอิง *",
                                       value=pf.get("_stock_code", ""),
                                       placeholder="เช่น 2383 HK")

        company_name = st.text_input("ชื่อบริษัทอ้างอิง (ภาษาอังกฤษ) *",
                                     value=pf.get("_company_name", ""),
                                     placeholder="เช่น Sunny Optical Technology (Group) Co., Ltd.")

        col5, col6 = st.columns(2)
        with col5:
            total_units = st.number_input("จำนวนหน่วยที่อนุมัติ *",
                                          min_value=1,
                                          value=pf.get("_total_units", 10_000_000_000),
                                          step=1_000_000_000,
                                          format="%d")
        with col6:
            ratio = st.number_input("อัตราอ้างอิง (1 : X) *",
                                    min_value=1,
                                    value=pf.get("_ratio", 100),
                                    step=1)

        # ── DATE PICKER ──
        default_date = pf.get("_trading_date", date.today())
        trading_date = st.date_input("ประมาณการวันซื้อขายวันแรก *",
                                     value=default_date,
                                     format="DD/MM/YYYY")

        filing_url = st.text_input("ลิงก์ข้อมูล Filing (สำหรับ QR Code)",
                                   value=pf.get("_filing_url", ""),
                                   placeholder="https://capital.sec.or.th/...")

        st.markdown("")
        btn_label = "💾 อัปเดตรายการ" if editing else "⚡ Generate Factsheet"
        submitted = st.form_submit_button(btn_label, use_container_width=True, type="primary")

    if submitted:
        # exchange_full and exchange_short_input are set above, outside the form
        current_exchange_full  = exchange_full
        current_exchange_short = exchange_short_input.strip()

        errors = []
        if not ticker.strip():               errors.append("กรุณากรอก Ticker")
        if not company_name.strip():         errors.append("กรุณากรอกชื่อบริษัท")
        if not stock_code.strip():           errors.append("กรุณากรอกรหัสหลักทรัพย์")
        if not current_exchange_short:       errors.append("กรุณากรอกชื่อย่อตลาด")

        if errors:
            for e in errors:
                st.error(e)
        else:
            data = build_data(
                ticker.strip(), company_name.strip(), stock_code.strip(),
                current_exchange_full, total_units, ratio,
                trading_date, filing_url.strip(), current_exchange_short
            )
            if editing:
                st.session_state.history[st.session_state.edit_index] = data
                save_history_to_disk(st.session_state.history)
                st.session_state.edit_index = None
                st.session_state.prefill    = {}
            else:
                st.session_state.history.append(data)

            save_history_to_disk(st.session_state.history)
            st.session_state.form_data = data
            st.rerun()

    if editing:
        if st.button("✖ ยกเลิกการแก้ไข", use_container_width=True):
            st.session_state.edit_index = None
            st.session_state.prefill    = {}
            st.rerun()

# ════════════════════
#  MIDDLE — PREVIEW
# ════════════════════
with preview_col:
    data = st.session_state.form_data

    if data is None:
        st.markdown("""
        <div style="background:#141720;border:1px dashed #252A3A;border-radius:12px;
                    padding:50px 20px;text-align:center;color:#5A637A;">
            <div style="font-size:36px;margin-bottom:10px;">📋</div>
            <div style="font-size:13px;font-weight:500;">กรอกข้อมูลด้านซ้ายแล้วกด Generate</div>
        </div>
        """, unsafe_allow_html=True)
    else:
        st.markdown('<div class="section-label">✅ ข้อมูลในตราสาร</div>', unsafe_allow_html=True)

        fields = [
            ("ชื่อย่อ",           data["ticker"]),
            ("ชื่อเต็ม (ไทย)",    data["full_name_thai"]),
            ("ชื่อเต็ม (อังกฤษ)", data["full_name_eng"]),
            ("หลักทรัพย์อ้างอิง", data["underlying_stock"]),
            ("ตลาดจดทะเบียน",     data["underlying_exchange"]),
            ("ชื่อย่อตลาด",       data["foreign_exchange"]),
            ("จำนวนหน่วย",        data["total_units"]),
            ("อัตราอ้างอิง",      data["ratio"]),
            ("วันซื้อขายวันแรก",  data["first_trading_date"]),
            ("ลิงก์ Filing",      data["filing_url"] or "—"),
        ]

        for label, value in fields:
            short_val = value if len(value) < 75 else value[:72] + "..."
            st.markdown(f"""
            <div class="gen-card">
                <div class="gen-label">{label}</div>
                <div class="gen-value">{short_val}</div>
            </div>
            """, unsafe_allow_html=True)

        st.markdown("---")
        st.markdown('<div class="section-label">⬇️ ดาวน์โหลด</div>', unsafe_allow_html=True)

        if st.session_state.template_bytes is None:
            st.warning("⚠️ กรุณาอัปโหลด template .pptx ก่อน")
        else:
            try:
                pptx_bytes = generate_pptx(st.session_state.template_bytes, data)
                st.download_button(
                    label=f"⬇️  ดาวน์โหลด {data['ticker']}.pptx",
                    data=pptx_bytes,
                    file_name=f"{data['ticker']}.pptx",
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                    use_container_width=True,
                    type="primary",
                )
            except Exception as e:
                st.error(f"❌ {str(e)}")

        if st.button("🗑️ ล้าง / สร้างใหม่", use_container_width=True):
            st.session_state.form_data  = None
            st.session_state.edit_index = None
            st.session_state.prefill    = {}
            st.rerun()

# ══════════════════════
#  RIGHT — HISTORY LOG
# ══════════════════════
with history_col:
    st.markdown('<div class="section-label">🕘 ประวัติการสร้าง</div>', unsafe_allow_html=True)

    if not st.session_state.history:
        st.markdown("""
        <div style="background:#141720;border:1px dashed #252A3A;border-radius:10px;
                    padding:30px;text-align:center;color:#5A637A;font-size:12px;">
            ยังไม่มีประวัติ<br>สร้าง DR แรกของคุณได้เลย
        </div>
        """, unsafe_allow_html=True)
    else:
        st.caption(f"{len(st.session_state.history)} รายการ")

        for i, h in enumerate(reversed(st.session_state.history)):
            real_idx = len(st.session_state.history) - 1 - i

            st.markdown(f"""
            <div class="hist-card">
                <span class="hist-ticker">{h['ticker']}</span>
                <div style="font-size:12px;color:#C8D0DC;margin-top:2px;">
                    {h['underlying_exchange']} · {h['ratio']}
                </div>
                <div class="hist-meta">{h['first_trading_date']} · {h.get('_created','')}</div>
            </div>
            """, unsafe_allow_html=True)

            b1, b2, b3 = st.columns(3)

            with b1:
                if st.button("👁 ดู", key=f"view_{real_idx}", use_container_width=True):
                    st.session_state.form_data  = h
                    st.session_state.edit_index = None
                    st.session_state.prefill    = {}
                    st.rerun()

            with b2:
                if st.button("✏️ แก้ไข", key=f"edit_{real_idx}", use_container_width=True):
                    st.session_state.edit_index = real_idx
                    st.session_state.prefill    = h
                    st.session_state.form_data  = None
                    st.rerun()

            with b3:
                if st.session_state.template_bytes:
                    try:
                        pptx_bytes = generate_pptx(st.session_state.template_bytes, h)
                        st.download_button(
                            label="⬇️",
                            data=pptx_bytes,
                            file_name=f"{h['ticker']}.pptx",
                            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                            key=f"dl_{real_idx}",
                            use_container_width=True,
                        )
                    except Exception:
                        st.button("⬇️", key=f"dl_err_{real_idx}",
                                  disabled=True, use_container_width=True)
                else:
                    st.button("⬇️", key=f"dl_none_{real_idx}",
                              disabled=True, help="อัปโหลด template ก่อน",
                              use_container_width=True)

            st.markdown("<div style='height:4px'></div>", unsafe_allow_html=True)

        st.markdown("---")
        if st.button("🗑️ ล้างประวัติทั้งหมด", use_container_width=True):
            st.session_state.history    = []
            save_history_to_disk([])
            st.session_state.form_data  = None
            st.session_state.edit_index = None
            st.session_state.prefill    = {}
            st.rerun()

# ── FOOTER ────────────────────────────────────────────────────
st.markdown("---")
st.markdown("""
<div style="text-align:center;font-size:11px;color:#5A637A;padding:6px 0;">
    DR Factsheet Generator · KTB Securities · ข้อมูลทั้งหมดเป็นความลับ
</div>
""", unsafe_allow_html=True)
